"""Generate the AE-TPF interview deck as a .pptx.

Run:  python scripts/build_deck.py
Out:  docs/slides/Zero-One-Games-AE-TPF.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- palette
BG     = RGBColor(0x0B, 0x10, 0x20)
PANEL  = RGBColor(0x12, 0x1A, 0x33)
CARD   = RGBColor(0x16, 0x20, 0x3F)
INK    = RGBColor(0xE8, 0xEC, 0xF8)
MUTED  = RGBColor(0x9F, 0xB0, 0xD6)
FAINT  = RGBColor(0x6B, 0x7A, 0xA3)
ACCENT = RGBColor(0x5B, 0x8C, 0xFF)
TEAL   = RGBColor(0x26, 0xD0, 0xA8)
AMBER  = RGBColor(0xFF, 0xB5, 0x47)
PINK   = RGBColor(0xFF, 0x6B, 0x8B)
LINE   = RGBColor(0x26, 0x33, 0x5C)

BODY_FONT = "Segoe UI"
MONO_FONT = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def _no_shadow(shape):
    el = shape._element.spPr
    # remove inherited shadow noise
    shape.shadow.inherit = False


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    _no_shadow(sp)
    return sp


def text(slide, x, y, w, h, runs, size=18, color=INK, bold=False, font=BODY_FONT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0, space_after=6):
    """runs: str OR list of (str, dict-overrides) OR list of paragraphs where a
    paragraph is a list of run-tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    if isinstance(runs, str):
        paragraphs = [[(runs, {})]]
    elif runs and isinstance(runs[0], tuple):
        paragraphs = [runs]
    else:
        paragraphs = runs

    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        if isinstance(para, tuple):
            para = [para]
        for run_text, ov in para:
            r = p.add_run()
            r.text = run_text
            r.font.size = Pt(ov.get("size", size))
            r.font.bold = ov.get("bold", bold)
            r.font.name = ov.get("font", font)
            r.font.color.rgb = ov.get("color", color)
    return tb


def base(eyebrow=None, title=None, eyebrow_color=TEAL, footer=None):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, SW, SH, fill=BG)
    # top accent hairline
    rect(slide, 0, 0, SW, Pt(6), fill=ACCENT)
    if eyebrow:
        text(slide, Inches(0.85), Inches(0.55), Inches(11), Inches(0.4),
             eyebrow.upper(), size=13, color=eyebrow_color, bold=True)
    if title:
        text(slide, Inches(0.82), Inches(0.95), Inches(11.6), Inches(1.2),
             title, size=34, color=INK, bold=True, line_spacing=1.0)
    if footer:
        rect(slide, Inches(0.85), Inches(7.02), Inches(11.6), Pt(1), fill=LINE)
        text(slide, Inches(0.85), Inches(7.08), Inches(11.6), Inches(0.35),
             footer, size=11, color=FAINT)
    return slide


def card(slide, x, y, w, h, heading, body, accent=ACCENT, hsize=15, bsize=12.5):
    rect(slide, x, y, w, h, fill=CARD, line=LINE, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(slide, x, y + Inches(0.18), Pt(4), h - Inches(0.36), fill=accent)
    tb = text(slide, x + Inches(0.28), y + Inches(0.18), w - Inches(0.5), h - Inches(0.34),
              [[(heading, {"size": hsize, "bold": True, "color": INK})],
               [(body, {"size": bsize, "color": MUTED})]],
              line_spacing=1.05, space_after=5)
    return tb


def chip(slide, x, y, label, color=TEAL, w=Inches(2.0)):
    h = Inches(0.42)
    rect(slide, x, y, w, h, fill=PANEL, line=LINE, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(slide, x, y + Inches(0.015), w, h, label, size=11.5, color=color,
         bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=MONO_FONT)


def bullets(slide, x, y, w, h, items, size=15, gap=8):
    paras = []
    for it in items:
        if isinstance(it, tuple):
            lead, rest = it
            paras.append([("▪  ", {"color": ACCENT, "bold": True, "size": size}),
                          (lead, {"color": INK, "bold": True, "size": size}),
                          (rest, {"color": MUTED, "size": size})])
        else:
            paras.append([("▪  ", {"color": ACCENT, "bold": True, "size": size}),
                          (it, {"color": MUTED, "size": size})])
    text(slide, x, y, w, h, paras, line_spacing=1.05, space_after=gap)


def codebox(slide, x, y, w, h, lines):
    rect(slide, x, y, w, h, fill=RGBColor(0x07, 0x0C, 0x1C), line=LINE, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    paras = [[(ln, {"font": MONO_FONT, "size": 12, "color": RGBColor(0xCF, 0xE0, 0xFF)})]
             for ln in lines]
    text(slide, x + Inches(0.25), y + Inches(0.18), w - Inches(0.45), h - Inches(0.3),
         paras, line_spacing=1.12, space_after=2)


# ================================================================ 1 TITLE
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=BG)
rect(s, 0, 0, SW, SH, fill=None)
# decorative bands
rect(s, 0, Inches(2.55), SW, Pt(3), fill=ACCENT)
rect(s, 0, Inches(2.6), Inches(4.4), Pt(3), fill=TEAL)
text(s, Inches(0.9), Inches(1.0), Inches(11.5), Inches(0.5),
     "ANALYTICS ENGINEERING · TPF TECHNICAL ROUND", size=15, color=TEAL, bold=True)
text(s, Inches(0.88), Inches(1.55), Inches(11.6), Inches(1.7),
     [[("Chef Master — Mobile Game", {"size": 46, "bold": True, "color": INK})],
      [("Analytics Platform", {"size": 46, "bold": True, "color": ACCENT})]],
     line_spacing=1.0)
text(s, Inches(0.9), Inches(3.0), Inches(10.8), Inches(1.0),
     [[("An end-to-end ", {"size": 19, "color": MUTED}),
       ("dbt + BigQuery + Semantic/BI", {"size": 19, "color": INK, "bold": True}),
       (" platform answering one question:", {"size": 19, "color": MUTED})],
      [("“Are we acquiring players profitably?”  →  D30 ROAS / LTV:CAC",
        {"size": 19, "color": INK, "bold": True})]],
     line_spacing=1.15, space_after=6)
chips = ["dbt Core/Fusion", "BigQuery", "MetricFlow SL", "GitHub CI/CD", "PII hashing"]
cx = Inches(0.9)
for c in chips:
    chip(s, cx, Inches(4.4), c, color=AMBER, w=Inches(2.05))
    cx += Inches(2.2)
text(s, Inches(0.9), Inches(6.6), Inches(11.6), Inches(0.5),
     "Sources unified: Firebase/GA4 · Adjust · Meta Ads · Unity Analytics",
     size=12.5, color=FAINT, font=MONO_FONT)

# ================================================================ 2 AGENDA
s = base("Agenda", "What I'll walk through", footer="≈90 min · 5 min intro · 15 min per pillar · 10 min Q&A")
items = [
    ("1 · Data Modeling & Architecture", "  — medallion layers + dimensional marts + ERD"),
    ("2 · Data Transformation", "  — incremental models, macros, dedup"),
    ("3 · Data Governance", "  — quality tests + PII security"),
    ("4 · Data Insight", "  — semantic layer, BI, the headline metric"),
    ("5 · DataOps", "  — branching, CI/CD, releases"),
]
bullets(s, Inches(0.9), Inches(2.05), Inches(7.4), Inches(4.5), items, size=18, gap=14)
card(s, Inches(8.7), Inches(2.05), Inches(3.75), Inches(2.05),
     "🎯 The business problem",
     "One quantifiable metric: D30 ROAS. ≥ 1 ⇒ a campaign pays back. Everything "
     "else is a driver of, or context for, that number.", accent=TEAL)
card(s, Inches(8.7), Inches(4.35), Inches(3.75), Inches(2.2),
     "🎁 Bonus section",
     "Anticipated interviewer Q&A with model answers — modeling, transformation, "
     "governance, insight, DataOps, and live “what-if” extensions.", accent=AMBER)

# ================================================================ 3 BUSINESS
s = base("The business problem", "One meaningful, quantifiable metric",
         footer="Driver metrics: eCPI · CTR/CVR · ARPDAU · ARPPU · retention · cohort LTV")
card(s, Inches(0.9), Inches(2.0), Inches(3.5), Inches(1.6), "D30 ROAS",
     "Cohort revenue ÷ ad spend. ≥ 1 ⇒ the campaign pays back inside the window.",
     accent=TEAL, hsize=22)
card(s, Inches(4.6), Inches(2.0), Inches(3.5), Inches(1.6), "LTV : CAC",
     "Long-horizon profitability the early ROAS window can't yet see.",
     accent=ACCENT, hsize=22)
card(s, Inches(8.3), Inches(2.0), Inches(4.15), Inches(1.6), "Why it matters",
     "Directly drives the UA budget: scale winners, cut losers, by "
     "network / campaign / geo.", accent=AMBER, hsize=18)
codebox(s, Inches(0.9), Inches(3.85), Inches(11.55), Inches(2.7), [
    "-- The headline query  (marts.ua_performance_daily)",
    "SELECT network,",
    "       SUM(spend)                                              AS spend,",
    "       SAFE_DIVIDE(SUM(cohort_revenue),",
    "                   NULLIF(SUM(spend), 0))                      AS roas",
    "FROM   marts.ua_performance_daily",
    "WHERE  date >= DATE_SUB(CURRENT_DATE(), INTERVAL 14 DAY)",
    "GROUP BY network ORDER BY spend DESC;",
])

# ================================================================ 4 ARCHITECTURE
s = base("Pillar 1 · Architecture", "Medallion layers + dimensional marts",
         footer="ERD + DAG: docs/architecture.md · mart reference: models/marts/README.md")
flow = [("Firebase · Adjust · Meta · Unity", FAINT), ("BRONZE\nraw", PINK),
        ("SILVER\nstaging views", ACCENT), ("GOLD\nincremental marts", TEAL),
        ("Semantic Layer\n+ BI", AMBER)]
fx = Inches(0.9)
fw = Inches(2.15)
for i, (lbl, col) in enumerate(flow):
    rect(s, fx, Inches(2.0), fw, Inches(0.95), fill=CARD, line=col, line_w=1.5,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, fx, Inches(2.0), fw, Inches(0.95), lbl, size=12, color=INK, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=MONO_FONT)
    fx += fw
    if i < len(flow) - 1:
        text(s, fx, Inches(2.0), Inches(0.18), Inches(0.95), "▶", size=14,
             color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        fx += Inches(0.18)
card(s, Inches(0.9), Inches(3.35), Inches(3.75), Inches(2.9), "🥈 Silver · staging",
     "One model per source object: type-cast, snake_case, flatten GA4 structs, "
     "hash PII, light dedup. Cheap views, always fresh. No business logic.",
     accent=ACCENT)
card(s, Inches(4.85), Inches(3.35), Inches(3.75), Inches(2.9), "🥇 Gold · marts",
     "9 business facts at a documented daily grain. Incremental tables, "
     "partitioned by date + clustered on hot keys. Grain-tested. The only "
     "layer BI touches.", accent=TEAL)
card(s, Inches(8.8), Inches(3.35), Inches(3.65), Inches(2.9), "🔭 Analytics",
     "Deeper player-behaviour exploration on the unified stg_events table.\n\n"
     "Conformed dims: date · app · platform · geo · network · app_version · "
     "spender_segment.", accent=AMBER)

# ================================================================ 5 PILLAR 1
s = base("Pillar 1 / 5", "Data Modeling & Architecture",
         eyebrow_color=ACCENT,
         footer="Evidence: docs/architecture.md · models/marts/README.md · marts.yml")
g = [("🏛️ Layered & dimensional", "Bronze→Silver→Gold; Kimball facts with conformed "
      "dimensions so ROAS/ARPDAU/LTV slice identically across every mart.", ACCENT),
     ("📏 Long-format facts", "spender_segment, step_num, days_since_install are ROWS — "
      "new segments/checkpoints need no schema change.", TEAL),
     ("⚙️ Tuned storage", "partition_by=date, cluster_by on hot keys, "
      "on_schema_change='append_new_columns'.", AMBER),
     ("📚 Documented", "Grain, source-of-truth precedence & caveats per mart; ERD + DAG; "
      "persist_docs pushes descriptions to BigQuery.", PINK)]
xs = [Inches(0.9), Inches(6.75)]
ys = [Inches(2.1), Inches(4.35)]
for i, (h, b, a) in enumerate(g):
    card(s, xs[i % 2], ys[i // 2], Inches(5.7), Inches(2.05), h, b, accent=a)

# ================================================================ 6 PILLAR 2
s = base("Pillar 2 / 5", "Data Transformation", eyebrow_color=ACCENT,
         footer="Live: dbt build -s +iap_arpdau_daily · macros/ · models/staging/firebase/")
g = [("🔁 Incremental everywhere", "6 marts + staging on insert_overwrite with 14/44-day "
      "windows — idempotent, absorbs late-arriving Adjust/Firebase data.", ACCENT),
     ("🧩 Macros / DRY", "ga4_flatten_event_columns() removes ~120 duplicated lines; "
      "iap.sql = one revenue definition; meta_get_action().", TEAL),
     ("🪟 Window-function dedup", "stg_events forward-fills level_id & derives prev/next "
      "level via run-numbering; cohorts via row_number().", AMBER),
     ("📦 Packages", "dbt_utils for surrogate keys, range & grain tests; codegen; "
      "audit_helper for refactor parity.", PINK)]
for i, (h, b, a) in enumerate(g):
    card(s, xs[i % 2], ys[i // 2], Inches(5.7), Inches(2.05), h, b, accent=a)

# ================================================================ 7 PILLAR 3
s = base("Pillar 3 / 5", "Data Governance — Quality & Security", eyebrow_color=ACCENT,
         footer="Found & fixed: retention_cohorts summed coin `amount` as USD → fixed via one iap_value_usd() definition")
g = [("✅ Quality, layered", "Grain unique_combination, ranges, accepted_values, source "
      "freshness, singular + unit tests (logic checks, no warehouse).", ACCENT),
     ("🔐 PII pseudonymisation", "Salted SHA-256 on user_id in staging — raw id never "
      "persisted, deterministic so joins still work.", TEAL),
     ("🛡️ Warehouse controls", "BigQuery policy tags + dynamic masking, authorized "
      "views, grants-as-code, row-level security design.", AMBER),
     ("🔑 Secrets", "All connection details via env_var(); keyless OIDC in CI; "
      "*.json/.bkp/.env git-ignored.", PINK)]
for i, (h, b, a) in enumerate(g):
    card(s, xs[i % 2], ys[i // 2], Inches(5.7), Inches(2.05), h, b, accent=a)

# ================================================================ 8 PILLAR 4
s = base("Pillar 4 / 5", "Data Insight", eyebrow_color=ACCENT,
         footer="Evidence: models/semantic/ · models/exposures.yml · docs/metrics_dictionary.md")
g = [("📊 Semantic Layer", "MetricFlow metrics: roas, ecpi, ctr/cvr, arpdau — defined "
      "once as numerator/denominator. Never “AVG a rate”.", ACCENT),
     ("🖥️ BI exposures", "4 dashboards wired into the dbt DAG (UA/ROAS, LTV & Retention, "
      "Creative, Engagement) — lineage warehouse → BI.", TEAL),
     ("📖 Metrics dictionary", "Every metric: definition, SQL logic, business impact — and "
      "why Adjust/Unity/Meta scopes don't reconcile.", AMBER),
     ("🧮 Recipe book", "Ready queries: ROAS, cohort LTV curves, retention, creative "
      "leaderboard, funnels.", PINK)]
for i, (h, b, a) in enumerate(g):
    card(s, xs[i % 2], ys[i // 2], Inches(5.7), Inches(2.05), h, b, accent=a)

# ================================================================ 9 PILLAR 5
s = base("Pillar 5 / 5", "DataOps", eyebrow_color=ACCENT,
         footer="Evidence: .github/workflows/ci.yml + cd.yml · .sqlfluff · .pre-commit-config.yaml · docs/dataops.md")
g = [("🌿 Trunk-based", "Protected main; short-lived branches; PR + green CI + CODEOWNERS "
      "review; squash-merge; tagged releases.", ACCENT),
     ("🤖 CI on every PR", "OIDC auth → dbt parse → sqlfluff → Slim CI build into a "
      "throwaway dataset → tests → drop dataset.", TEAL),
     ("🚀 Gated CD", "Merge to main → freshness → dbt build --target prod (approval "
      "gate) → publish manifest + docs.", AMBER),
     ("🧹 Pre-commit", "sqlfluff, secret detection, dbt-checkpoint (marts must have "
      "description + tests).", PINK)]
for i, (h, b, a) in enumerate(g):
    card(s, xs[i % 2], ys[i // 2], Inches(5.7), Inches(2.05), h, b, accent=a)

# ================================================================ 10 SIGNALS
s = base("Senior signals", "What I'd point at first",
         footer="Both vertical depth and horizontal breadth — balanced across the 5 pillars")
bullets(s, Inches(0.9), Inches(2.05), Inches(11.6), Inches(4.6), [
    ("Source-of-truth design:", " Adjust vs Unity vs Meta measure different scopes — "
     "modelled deliberately, labelled, never silently mixed."),
    ("Correctness judgement:", " caught the IAP `amount` inflation bug; centralised the "
     "fix so every mart reconciles on one definition."),
    ("Idempotent incrementals:", " insert_overwrite + lookback windows handle retroactive "
     "attribution without duplicates."),
    ("Governance as code:", " PII hashing, policy-tag design, grants, secrets in env — "
     "all version-controlled and reviewed in PRs."),
    ("Cheap, safe CI:", " Slim CI builds only what changed into a disposable dataset; "
     "production is never touched by a pull request."),
], size=17, gap=16)

# ================================================================ 11 ROADMAP
s = base("Roadmap", "What I'd build next",
         footer="Honest about current limits — and how I'd close them")
g = [("📈 IAP + ad cohort LTV", "Ingest the Adjust Cohorts API to add ad-revenue LTV and "
      "true d28 ROAS (Firebase ad events carry no $ value today).", TEAL),
     ("👥 Exact DAU", "Re-pull Unity without clientVersion in groupBy to remove the "
      "DAU double-count upper bound.", ACCENT),
     ("🧊 Snapshots / SCD2", "dbt snapshots on slowly-changing campaign metadata for "
      "point-in-time attribution.", AMBER),
     ("🔔 Anomaly alerts", "Elementary / freshness + volume anomaly tests piped to Slack "
      "on the daily run.", PINK)]
for i, (h, b, a) in enumerate(g):
    card(s, xs[i % 2], ys[i // 2], Inches(5.7), Inches(2.05), h, b, accent=a)

# ================================================================ 12 DEMO PLAN
s = base("Demo plan", "Time-boxed walkthrough (Show → Why → What-if)",
         footer="Full script: docs/demo_plan.md · keep answers concise, breadth = depth")
rows = [
    ("00–05", "Intro", "Business problem + the metric + repo map", TEAL),
    ("05–20", "Modeling", "DAG/ERD, layers, one mart's grain", ACCENT),
    ("20–35", "Transformation", "dbt build a mart; macros; incrementals", ACCENT),
    ("35–50", "Governance", "tests run; PII hashing; secrets", ACCENT),
    ("50–65", "Insight", "semantic metrics; exposures; ROAS query", ACCENT),
    ("65–80", "DataOps", "CI/CD walkthrough; branching; rollback", ACCENT),
    ("80–90", "Q&A", "caveats, scope mismatch, the bug I fixed", AMBER),
]
y = Inches(2.05)
for mins, topic, what, col in rows:
    rect(s, Inches(0.9), y, Inches(1.35), Inches(0.6), fill=PANEL, line=LINE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, Inches(0.9), y, Inches(1.35), Inches(0.6), mins, size=12.5, color=col,
         bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=MONO_FONT)
    text(s, Inches(2.45), y, Inches(2.6), Inches(0.6), topic, size=15, color=INK,
         bold=True, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(5.0), y, Inches(7.4), Inches(0.6), what, size=13.5, color=MUTED,
         anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.66)

# ================================================================ 13 DIVIDER
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=PANEL)
rect(s, 0, Inches(3.05), SW, Pt(4), fill=AMBER)
rect(s, 0, Inches(3.1), Inches(4.0), Pt(4), fill=TEAL)
text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.5),
     "BONUS", size=16, color=AMBER, bold=True)
text(s, Inches(0.88), Inches(2.45), Inches(11.6), Inches(1.0),
     "Anticipated Q&A — with solutions", size=40, color=INK, bold=True)
text(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(0.6),
     "The questions I expect, and how I'd answer them under the Show → Why → What-if frame.",
     size=17, color=MUTED)

# ================================================================ Q&A slides
def qa_slide(eyebrow, title, pairs, footer=None):
    s = base(eyebrow, title, eyebrow_color=AMBER, footer=footer)
    y = Inches(2.0)
    h = Inches(2.18) if len(pairs) == 2 else Inches(1.42)
    for q, a in pairs:
        rect(s, Inches(0.9), y, Inches(11.55), h, fill=CARD, line=LINE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, Inches(0.9), y + Inches(0.16), Pt(4), h - Inches(0.32), fill=AMBER)
        text(s, Inches(1.2), y + Inches(0.16), Inches(11.0), h - Inches(0.3),
             [[("Q.  ", {"color": AMBER, "bold": True, "size": 15}),
               (q, {"color": INK, "bold": True, "size": 15})],
              [("A.  ", {"color": TEAL, "bold": True, "size": 13.5}),
               (a, {"color": MUTED, "size": 13.5})]],
             line_spacing=1.05, space_after=6)
        y += h + Inches(0.25)
    return s


qa_slide("Bonus Q&A · Modeling", "Data Modeling & Architecture", [
    ("Why medallion + Kimball, not Data Vault or 3NF?",
     "Single small team, analytics-first, fast-moving sources. Medallion gives clear "
     "testable layers; dimensional facts on conformed dims keep BI simple. Data Vault's "
     "hubs/links/sats add ingestion overhead we don't need at this scale."),
    ("Why long-format facts instead of wide columns?",
     "spender_segment, step_num and days_since_install are dimensions, not columns. As "
     "rows, adding D45 LTV or a new segment is a config bump (max_dsi) — no schema change, "
     "no BI rework, and aggregations stay summable."),
])

qa_slide("Bonus Q&A · Modeling", "Architecture (cont.)", [
    ("Why are staging views but marts incremental tables?",
     "Staging is light passthrough — views cost nothing and are always fresh. Marts back "
     "dashboards and scan history, so we materialise them, partition by date and cluster "
     "on hot keys to cut cost and latency."),
    ("How do you keep metrics consistent across marts?",
     "Conformed dimensions + shared macros (iap_value_usd) + a semantic layer. One "
     "definition, referenced everywhere, so ROAS means the same thing in every mart and "
     "every dashboard."),
])

qa_slide("Bonus Q&A · Transformation", "Data Transformation", [
    ("Why insert_overwrite over merge for incrementals?",
     "Adjust/Firebase mutate retroactively. insert_overwrite atomically replaces whole "
     "date partitions, so re-running the 14/44-day window is idempotent and dupe-free — "
     "cheaper than a row-level merge on BigQuery."),
    ("Why a macro for the GA4 flattening — what's the trade-off?",
     "The daily and intraday shards must project identical columns or the UNION fails. "
     "One macro guarantees that and removes ~120 duplicated lines. Trade-off: a layer of "
     "indirection — mitigated by clear naming + comments."),
])

qa_slide("Bonus Q&A · Governance & Security", "Quality & Security", [
    ("How do you protect PII but keep joins working?",
     "Deterministic salted SHA-256: same input → same hash, so joins on the hash still "
     "work, but the raw id is never stored. Salt comes from env (DBT_PII_SALT), defeating "
     "rainbow tables."),
    ("Hashing vs dynamic data masking vs policy tags — when each?",
     "Hash when the raw value is never needed downstream (user_id). Policy tags + dynamic "
     "masking when privileged roles still need the real value (geo, pseudo_id). "
     "Authorized views + IAM control who reaches the raw layer at all."),
])

qa_slide("Bonus Q&A · Governance & Security", "Quality & Security (cont.)", [
    ("How do you keep CI green but still catch data drift?",
     "Severity strategy: hard invariants (grain uniqueness, null keys, future dates) are "
     "errors; noisy range expectations are warnings. Build stays green; drift still "
     "surfaces in the CI log and docs."),
    ("A revenue number looks wrong — how do you debug it?",
     "Exactly how I found the retention_cohorts bug: trace to the single definition. It "
     "summed the generic `amount` coin param as USD across all events. Fix was centralising "
     "iap_value_usd() so all marts reconcile — now covered by a singular test."),
])

qa_slide("Bonus Q&A · Data Insight", "Data Insight", [
    ("Why a semantic layer instead of metrics in the BI tool?",
     "Define once, consume everywhere (BI, notebooks, ad-hoc) with governed numerator/"
     "denominator so nobody AVG()s a rate. BI-only metrics drift between dashboards and "
     "can't be reused or tested."),
    ("Why don't Adjust, Unity and Meta ROAS reconcile — is that a bug?",
     "No — different scopes. Adjust = paid-UA attributed users; Unity = whole game; Meta = "
     "its own view-through attribution. They measure different populations; dashboards "
     "label them distinctly rather than forcing a false match."),
])

qa_slide("Bonus Q&A · DataOps", "DataOps", [
    ("How does a PR get validated without touching production?",
     "Slim CI: build only state:modified+ into a throwaway dataset (dbt_ci_pr_<n>), "
     "deferring unchanged refs to the prod manifest, run their tests, then drop the "
     "dataset. Prod is never written by a PR."),
    ("How do you roll back a bad deploy?",
     "Releases are tagged on main; redeploy the previous tag. Because marts are "
     "partition-replace incrementals, re-running a known-good commit restores state; "
     "--full-refresh rebuilds from source if needed."),
])

qa_slide("Bonus Q&A · Live “What-if”", "Extensions I'm ready to do live", [
    ("Add D45 / D60 to the LTV curve?",
     "Bump max_dsi in cohort_ltv_daily, widen the incremental window, dbt run "
     "--full-refresh. Long format means no schema change and the BI curve just extends."),
    ("Onboard a new ad network tomorrow?",
     "Add the raw table + src yml, a stg_<src>__* view, then fold into "
     "ua_performance_daily (keeps the source-of-truth pattern) — or a dedicated mart if "
     "the grain differs."),
], footer="Also ready: add a new metric to the semantic layer, write a singular test, add a masking policy — live.")

# ================================================================ LAST
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, fill=BG)
rect(s, 0, Inches(3.0), SW, Pt(4), fill=ACCENT)
rect(s, 0, Inches(3.05), Inches(4.0), Pt(4), fill=TEAL)
text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.5),
     "LET'S DIG IN", size=16, color=TEAL, bold=True)
text(s, Inches(0.88), Inches(2.4), Inches(11.6), Inches(1.0),
     "Live demo & Q&A", size=46, color=INK, bold=True)
text(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(0.8),
     "Repo is navigable · dbt parses clean · tests pass · semantic metrics & dashboards ready.",
     size=18, color=MUTED)
cx = Inches(0.9)
for c in ["README.md", "docs/demo_plan.md", "docs/architecture.md", "marts/README.md"]:
    chip(s, cx, Inches(4.7), c, color=AMBER, w=Inches(2.55))
    cx += Inches(2.75)
text(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.4),
     "Thank you — Analytics Engineering · TPF", size=13, color=FAINT, font=MONO_FONT)

# ---------------------------------------------------------------- save
out = Path(__file__).resolve().parent.parent / "docs" / "slides" / "Zero-One-Games-AE-TPF.pptx"
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(out))
print(f"Wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
